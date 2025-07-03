"""
Content Analysis API - Task 1.1
Intelligent auto-categorization for educational content using LLM analysis
"""

import json
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import tempfile
import os

try:
    from ollama import Client
    OLLAMA_AVAILABLE = True
except ImportError:
    OLLAMA_AVAILABLE = False
    logging.warning("Ollama not available - content analysis will use fallback only")

# Content extraction libraries
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Educational content categories based on current system
EDUCATIONAL_CATEGORIES = [
    "lesson-plan", "worksheet", "assessment", "resource", "activity"
]

SUBJECT_AREAS = [
    "English", "Religious Education", "Learning Support", "Other"
]

DIFFICULTY_LEVELS = [
    "beginner", "intermediate", "advanced"
]

GRADE_TARGETS = [
    "early-years", "primary", "secondary", "adult-ed"
]

def extract_text_from_file(file_path: str, mime_type: str = None) -> str:
    """Extract text content from various file formats"""
    try:
        file_path = Path(file_path)
        file_extension = file_path.suffix.lower()
        
        # Text files
        if file_extension in ['.txt', '.md', '.rst']:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        # PDF files
        if file_extension == '.pdf' and PDF_AVAILABLE:
            try:
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                    return text
            except Exception as e:
                logging.warning(f"PDF extraction failed: {e}")
        
        # Word documents
        if file_extension in ['.docx', '.doc'] and DOCX_AVAILABLE:
            try:
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            except Exception as e:
                logging.warning(f"DOCX extraction failed: {e}")
        
        # PowerPoint presentations
        if file_extension in ['.pptx', '.ppt'] and PPTX_AVAILABLE:
            try:
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except Exception as e:
                logging.warning(f"PPTX extraction failed: {e}")
        
        # Fallback - try to read as text
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception:
            pass
            
        return ""
        
    except Exception as e:
        logging.error(f"Content extraction failed for {file_path}: {e}")
        return ""

def analyze_educational_content(title: str, content: str, filename: str = "", 
                              client: Client = None, model: str = "qwen2.5:7b") -> Dict[str, Any]:
    """
    Analyze educational content using LLM for intelligent categorization
    Adapted from MCP-Testing/smart_tagging_bridge.py for educational context
    """
    
    if not client or not OLLAMA_AVAILABLE:
        return _fallback_analysis(title, content, filename)
    
    # Educational-specific prompt engineering
    analysis_prompt = f"""Analyze the following educational content and categorize it appropriately.

Title: "{title}"
Filename: "{filename}"
Content: "{content[:1000]}{'...' if len(content) > 1000 else ''}"

Please categorize this content using these categories:

CONTENT TYPES: {', '.join(EDUCATIONAL_CATEGORIES)}
SUBJECTS: {', '.join(SUBJECT_AREAS)}
DIFFICULTY: {', '.join(DIFFICULTY_LEVELS)}
GRADE LEVEL: {', '.join(GRADE_TARGETS)}

Instructions:
- Choose the MOST appropriate content type from: {', '.join(EDUCATIONAL_CATEGORIES)}
- Choose the MOST appropriate subject from: {', '.join(SUBJECT_AREAS)}
- Estimate difficulty level from: {', '.join(DIFFICULTY_LEVELS)}
- Estimate target grade level from: {', '.join(GRADE_TARGETS)}
- Assign a confidence score (0.0-1.0) for each classification
- Suggest 2-4 relevant tags for organization

Respond with ONLY a JSON object in this format:
{{
    "content_type": "lesson-plan",
    "content_type_confidence": 0.95,
    "subject": "English",
    "subject_confidence": 0.87,
    "difficulty": "intermediate",
    "difficulty_confidence": 0.78,
    "grade_level": "primary",
    "grade_level_confidence": 0.82,
    "suggested_tags": ["reading", "comprehension", "vocabulary"],
    "overall_confidence": 0.85
}}

Classification:"""

    try:
        response = client.chat(
            model=model,
            messages=[{
                "role": "user",
                "content": analysis_prompt
            }]
        )
        
        response_text = (response.message.content or "").strip()
        
        try:
            # Parse JSON response
            analysis = json.loads(response_text)
            
            # Validate response structure
            if isinstance(analysis, dict) and all(key in analysis for key in 
                ['content_type', 'subject', 'difficulty', 'grade_level']):
                
                # Ensure confidence scores
                for field in ['content_type', 'subject', 'difficulty', 'grade_level']:
                    confidence_key = f"{field}_confidence"
                    if confidence_key not in analysis:
                        analysis[confidence_key] = 0.7  # Default confidence
                
                if 'overall_confidence' not in analysis:
                    # Calculate overall confidence as average
                    confidences = [analysis.get(f"{field}_confidence", 0.7) 
                                 for field in ['content_type', 'subject', 'difficulty', 'grade_level']]
                    analysis['overall_confidence'] = sum(confidences) / len(confidences)
                
                if 'suggested_tags' not in analysis:
                    analysis['suggested_tags'] = []
                
                return analysis
        
        except json.JSONDecodeError:
            logging.warning("LLM response was not valid JSON, using fallback analysis")
    
    except Exception as e:
        logging.error(f"LLM analysis failed: {e}")
    
    # Fallback to keyword-based analysis
    return _fallback_analysis(title, content, filename)

def _fallback_analysis(title: str, content: str, filename: str = "") -> Dict[str, Any]:
    """
    Keyword-based fallback analysis for when LLM is unavailable
    Based on the fallback logic from MCP-Testing/smart_tagging_bridge.py
    """
    
    combined_text = (title + " " + content + " " + filename).lower()
    
    # Content type detection
    content_type = "resource"  # default
    content_type_confidence = 0.5
    
    if any(word in combined_text for word in ["lesson", "plan", "teaching", "instruction", "objective"]):
        content_type = "lesson-plan"
        content_type_confidence = 0.8
    elif any(word in combined_text for word in ["worksheet", "exercise", "practice", "activity", "handout"]):
        content_type = "worksheet"
        content_type_confidence = 0.8
    elif any(word in combined_text for word in ["test", "quiz", "exam", "assessment", "evaluation"]):
        content_type = "assessment"
        content_type_confidence = 0.8
    elif any(word in combined_text for word in ["activity", "game", "project", "experiment"]):
        content_type = "activity"
        content_type_confidence = 0.7
    
    # Subject detection
    subject = "Other"
    subject_confidence = 0.5
    
    if any(word in combined_text for word in ["english", "reading", "writing", "literature", "grammar", "spelling"]):
        subject = "English"
        subject_confidence = 0.8
    elif any(word in combined_text for word in ["math", "mathematics", "number", "calculation", "geometry", "algebra"]):
        subject = "Other"  # Mathematics maps to Other
        subject_confidence = 0.7
    elif any(word in combined_text for word in ["science", "biology", "chemistry", "physics", "experiment"]):
        subject = "Other"  # Science maps to Other
        subject_confidence = 0.7
    elif any(word in combined_text for word in ["religion", "religious", "faith", "prayer", "christian", "catholic"]):
        subject = "Religious Education"
        subject_confidence = 0.8
    elif any(word in combined_text for word in ["support", "special", "needs", "inclusion", "accessibility"]):
        subject = "Learning Support"
        subject_confidence = 0.8
    
    # Difficulty detection
    difficulty = "intermediate"
    difficulty_confidence = 0.6
    
    if any(word in combined_text for word in ["basic", "simple", "easy", "introduction", "beginner"]):
        difficulty = "beginner"
        difficulty_confidence = 0.7
    elif any(word in combined_text for word in ["advanced", "complex", "difficult", "challenging", "expert"]):
        difficulty = "advanced"
        difficulty_confidence = 0.7
    
    # Grade level detection
    grade_level = "primary"
    grade_level_confidence = 0.6
    
    if any(word in combined_text for word in ["nursery", "reception", "early", "preschool", "kindergarten"]):
        grade_level = "early-years"
        grade_level_confidence = 0.8
    elif any(word in combined_text for word in ["secondary", "high school", "gcse", "year 7", "year 8", "year 9", "year 10", "year 11"]):
        grade_level = "secondary"
        grade_level_confidence = 0.8
    elif any(word in combined_text for word in ["adult", "mature", "university", "college", "professional"]):
        grade_level = "adult-ed"
        grade_level_confidence = 0.7
    
    # Generate suggested tags
    suggested_tags = []
    if "reading" in combined_text:
        suggested_tags.append("reading")
    if "writing" in combined_text:
        suggested_tags.append("writing")
    if "number" in combined_text or "math" in combined_text:
        suggested_tags.append("numeracy")
    if "comprehension" in combined_text:
        suggested_tags.append("comprehension")
    
    # Limit to 4 tags
    suggested_tags = suggested_tags[:4]
    
    # Calculate overall confidence
    overall_confidence = (content_type_confidence + subject_confidence + 
                         difficulty_confidence + grade_level_confidence) / 4
    
    return {
        "content_type": content_type,
        "content_type_confidence": content_type_confidence,
        "subject": subject,
        "subject_confidence": subject_confidence,
        "difficulty": difficulty,
        "difficulty_confidence": difficulty_confidence,
        "grade_level": grade_level,
        "grade_level_confidence": grade_level_confidence,
        "suggested_tags": suggested_tags,
        "overall_confidence": overall_confidence,
        "analysis_method": "fallback"
    }

def analyze_content_endpoint(file, metadata: Dict[str, str]) -> Dict[str, Any]:
    """
    Main endpoint function for content analysis
    Handles file processing and calls analysis functions
    """
    
    # Extract metadata
    title = metadata.get('title', '')
    description = metadata.get('description', '')
    filename = metadata.get('filename', getattr(file, 'filename', ''))
    
    # Save file temporarily for content extraction
    temp_file = None
    try:
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, 
                                       suffix=Path(filename).suffix if filename else '.tmp') as tf:
            if hasattr(file, 'save'):
                file.save(tf.name)
            else:
                tf.write(file.read())
            temp_file = tf.name
        
        # Extract content from file
        content = extract_text_from_file(temp_file, getattr(file, 'content_type', None))
        
        # Combine all text for analysis
        combined_content = f"{title}\n{description}\n{content}".strip()
        
        # Initialize Ollama client if available
        client = None
        if OLLAMA_AVAILABLE:
            try:
                client = Client()
                # Test connection
                client.list()
            except Exception as e:
                logging.warning(f"Ollama connection failed: {e}")
                client = None
        
        # Perform analysis
        analysis = analyze_educational_content(
            title=title,
            content=combined_content,
            filename=filename,
            client=client
        )
        
        # Add metadata
        analysis['analysis_timestamp'] = str(Path().cwd())  # Use current time in production
        analysis['has_llm'] = client is not None
        analysis['extracted_content_length'] = len(content)
        
        return {
            'status': 'success',
            'analysis': analysis,
            'metadata': {
                'content_extracted': len(content) > 0,
                'analysis_method': 'llm' if client else 'fallback',
                'filename': filename
            }
        }
        
    except Exception as e:
        logging.error(f"Content analysis failed: {e}")
        return {
            'status': 'error',
            'message': str(e),
            'analysis': _fallback_analysis(title, description, filename)
        }
    
    finally:
        # Clean up temporary file
        if temp_file and os.path.exists(temp_file):
            try:
                os.remove(temp_file)
            except Exception as e:
                logging.warning(f"Failed to clean up temp file: {e}") 