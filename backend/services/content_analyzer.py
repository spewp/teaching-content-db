#!/usr/bin/env python3
"""
ContentAnalyzer Service - Task 2.1: Local LLM Integration
Educational content analysis using Ollama LLM following mcp-smart-notes pattern
Based on MCP-Testing/smart_tagging_bridge.py implementation
"""

import json
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import tempfile
import os
import re

try:
    from ollama import Client
    OLLAMA_AVAILABLE = True
except ImportError:
    OLLAMA_AVAILABLE = False
    logging.warning("Ollama not available - content analysis will use fallback only")

# Content extraction libraries
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not available - Word document support disabled")

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logging.warning("PyPDF2 not available - PDF support disabled")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available - PowerPoint support disabled")

# Educational content categories - Enhanced from mcp-smart-notes pattern
EDUCATIONAL_CATEGORIES = [
    "lesson-plan", "worksheet", "assessment", "resource", "activity"
]

SUBJECT_AREAS = [
    "English", "Religious Education", "Learning Support", "Other"
]

DIFFICULTY_LEVELS = [
    "beginner", "intermediate", "advanced"
]

GRADE_TARGETS = [
    "early-years", "primary", "secondary", "adult-ed"
]

class ContentAnalyzer:
    """
    Educational content analyzer using Local LLM integration
    Following mcp-smart-notes pattern from smart_tagging_bridge.py
    """
    
    def __init__(self, model: str = "qwen2.5:7b"):
        """
        Initialize ContentAnalyzer with Ollama client
        Following mcp-smart-notes initialization pattern
        """
        self.model = model
        self.client = None
        
        # Initialize Ollama client if available
        if OLLAMA_AVAILABLE:
            try:
                self.client = Client()
                # Test connection by listing models with timeout
                models = self.client.list()
                if models and hasattr(models, 'models') and len(models.models) > 0:
                    logging.info(f"✅ Ollama connected with {len(models.models)} models")
                    logging.info(f"🤖 Using model: {self.model}")
                    # Verify our specific model is available
                    model_names = [m.model for m in models.models]
                    if self.model not in model_names:
                        logging.warning(f"⚠️ Model {self.model} not found. Available: {model_names}")
                else:
                    logging.warning("⚠️ Ollama connected but no models available")
                    self.client = None
            except Exception as e:
                logging.warning(f"Ollama connection failed: {e}")
                self.client = None
        
        # Log initialization status
        if self.client:
            logging.info(f"🚀 ContentAnalyzer initialized with LLM support")
            logging.info(f"🏷️ Educational categories: {', '.join(EDUCATIONAL_CATEGORIES)}")
        else:
            logging.warning("⚠️ ContentAnalyzer initialized with fallback analysis only")
    
    def analyze_educational_content(self, title: str, content: str, filename: str = "") -> Dict[str, Any]:
        """
        Analyze educational content using LLM for intelligent categorization
        Based on analyze_content_for_tags from mcp-smart-notes but adapted for educational context
        LLM-ONLY MODE - No fallback analysis
        """
        
        if not self.client:
            raise Exception("LLM not available - Ollama client not connected. Please ensure Ollama is running with 'ollama serve'")
        
        # Educational-specific prompt engineering - Enhanced from mcp-smart-notes pattern
        analysis_prompt = f"""Analyze the following educational content and categorize it appropriately for a teaching database.

Title: "{title}"
Filename: "{filename}"
Content: "{content[:1500]}{'...' if len(content) > 1500 else ''}"

Please categorize this educational content using these specific categories:

CONTENT TYPES: {', '.join(EDUCATIONAL_CATEGORIES)}
- lesson-plan: Full lesson plans with objectives, activities, and outcomes
- worksheet: Practice exercises, handouts, and student activities
- assessment: Tests, quizzes, evaluations, and rubrics
- resource: Reference materials, guides, and supplementary content
- activity: Interactive activities, games, and projects

SUBJECTS: {', '.join(SUBJECT_AREAS)}
- English: Reading, writing, literature, grammar, spelling
- Religious Education: Faith-based content, prayer, religious studies
- Learning Support: Special needs, inclusion, accessibility materials
- Other: Cross-curricular or unspecified subjects

DIFFICULTY LEVELS: {', '.join(DIFFICULTY_LEVELS)}
- beginner: Basic concepts, simple vocabulary, foundational skills
- intermediate: Standard grade-level content, moderate complexity
- advanced: Complex concepts, challenging vocabulary, higher-order thinking

GRADE TARGETS: {', '.join(GRADE_TARGETS)}
- early-years: Ages 3-5, nursery, reception, kindergarten
- primary: Ages 5-11, elementary school, years 1-6
- secondary: Ages 11-18, high school, years 7-13
- adult-ed: Adult learners, professional development, continuing education

Instructions:
- Choose the MOST appropriate category from each list
- Only use categories from the exact lists provided above
- Assign a confidence score (0.0-1.0) for each classification
- Suggest 2-4 relevant organizational tags
- Consider the educational context and teaching purpose

Respond with ONLY a JSON object in this exact format:
{{
    "content_type": "lesson-plan",
    "content_type_confidence": 0.95,
    "subject": "English",
    "subject_confidence": 0.87,
    "difficulty": "intermediate",
    "difficulty_confidence": 0.78,
    "grade_level": "primary",
    "grade_level_confidence": 0.82,
    "suggested_tags": ["reading", "comprehension", "vocabulary"],
    "overall_confidence": 0.85
}}

Classification:"""

        try:
            # Call LLM using mcp-smart-notes pattern
            response = self.client.chat(
                model=self.model,
                messages=[{
                    "role": "user",
                    "content": analysis_prompt
                }]
            )
            
            # Extract response text safely - following mcp-smart-notes pattern
            response_text = (response.message.content or "").strip()
            
            # Try to parse as JSON - following mcp-smart-notes pattern
            try:
                analysis = json.loads(response_text)
                
                # Validate response structure
                if isinstance(analysis, dict) and self._validate_analysis_response(analysis):
                    
                    # Ensure all required fields and confidence scores
                    analysis = self._normalize_analysis_response(analysis)
                    
                    # Add metadata
                    analysis['analysis_method'] = 'llm'
                    analysis['model_used'] = self.model
                    
                    logging.info(f"✅ LLM analysis completed with {analysis['overall_confidence']:.2f} confidence")
                    return analysis
                else:
                    raise Exception("LLM returned invalid response structure - missing required fields")
                    
            except json.JSONDecodeError:
                raise Exception("LLM returned invalid JSON response - unable to parse analysis results")
                
        except Exception as e:
            logging.error(f"❌ LLM analysis failed: {e}")
            raise Exception(f"LLM analysis failed: {str(e)}")
    
    def _validate_analysis_response(self, analysis: Dict[str, Any]) -> bool:
        """Validate that LLM response has required fields"""
        required_fields = ['content_type', 'subject', 'difficulty', 'grade_level']
        
        # Check required fields exist
        if not all(field in analysis for field in required_fields):
            return False
        
        # Validate values are from allowed lists
        if analysis['content_type'] not in EDUCATIONAL_CATEGORIES:
            return False
        if analysis['subject'] not in SUBJECT_AREAS:
            return False
        if analysis['difficulty'] not in DIFFICULTY_LEVELS:
            return False
        if analysis['grade_level'] not in GRADE_TARGETS:
            return False
        
        return True
    
    def _normalize_analysis_response(self, analysis: Dict[str, Any]) -> Dict[str, Any]:
        """Normalize LLM response to ensure all required fields and confidence scores"""
        
        # Ensure confidence scores exist
        for field in ['content_type', 'subject', 'difficulty', 'grade_level']:
            confidence_key = f"{field}_confidence"
            if confidence_key not in analysis:
                analysis[confidence_key] = 0.7  # Default confidence
        
        # Calculate overall confidence if missing
        if 'overall_confidence' not in analysis:
            confidences = [analysis.get(f"{field}_confidence", 0.7) 
                         for field in ['content_type', 'subject', 'difficulty', 'grade_level']]
            analysis['overall_confidence'] = sum(confidences) / len(confidences)
        
        # Ensure suggested tags exist
        if 'suggested_tags' not in analysis:
            analysis['suggested_tags'] = []
        
        # Limit suggested tags to 4 max
        if len(analysis['suggested_tags']) > 4:
            analysis['suggested_tags'] = analysis['suggested_tags'][:4]
        
        return analysis
    
    def _fallback_analysis(self, title: str, content: str, filename: str = "") -> Dict[str, Any]:
        """
        Keyword-based fallback analysis for when LLM is unavailable
        Following mcp-smart-notes fallback pattern from smart_tagging_bridge.py
        """
        
        combined_text = (title + " " + content + " " + filename).lower()
        
        # Content type detection - following mcp-smart-notes keyword pattern
        content_type = "resource"  # default
        content_type_confidence = 0.5
        
        if any(word in combined_text for word in ["lesson", "plan", "teaching", "instruction", "objective", "learning goal"]):
            content_type = "lesson-plan"
            content_type_confidence = 0.8
        elif any(word in combined_text for word in ["worksheet", "exercise", "practice", "activity", "handout", "task"]):
            content_type = "worksheet"
            content_type_confidence = 0.8
        elif any(word in combined_text for word in ["test", "quiz", "exam", "assessment", "evaluation", "rubric"]):
            content_type = "assessment"
            content_type_confidence = 0.8
        elif any(word in combined_text for word in ["activity", "game", "project", "experiment", "investigation"]):
            content_type = "activity"
            content_type_confidence = 0.7
        
        # Subject detection - enhanced for educational context
        subject = "Other"
        subject_confidence = 0.5
        
        if any(word in combined_text for word in ["english", "reading", "writing", "literature", "grammar", "spelling", "phonics"]):
            subject = "English"
            subject_confidence = 0.8
        elif any(word in combined_text for word in ["math", "mathematics", "number", "calculation", "geometry", "algebra", "arithmetic"]):
            subject = "Mathematics"
            subject_confidence = 0.8
        elif any(word in combined_text for word in ["science", "biology", "chemistry", "physics", "experiment", "scientific"]):
            subject = "Science"
            subject_confidence = 0.8
        elif any(word in combined_text for word in ["religion", "religious", "faith", "prayer", "christian", "catholic", "bible"]):
            subject = "Religious Education"
            subject_confidence = 0.8
        elif any(word in combined_text for word in ["support", "special", "needs", "inclusion", "accessibility", "sen"]):
            subject = "Learning Support"
            subject_confidence = 0.8
        
        # Difficulty detection - enhanced pattern
        difficulty = "intermediate"
        difficulty_confidence = 0.6
        
        if any(word in combined_text for word in ["basic", "simple", "easy", "introduction", "beginner", "foundation"]):
            difficulty = "beginner"
            difficulty_confidence = 0.7
        elif any(word in combined_text for word in ["advanced", "complex", "difficult", "challenging", "expert", "higher order"]):
            difficulty = "advanced"
            difficulty_confidence = 0.7
        
        # Grade level detection - enhanced for educational context
        grade_level = "primary"
        grade_level_confidence = 0.6
        
        if any(word in combined_text for word in ["nursery", "reception", "early", "preschool", "kindergarten", "eyfs"]):
            grade_level = "early-years"
            grade_level_confidence = 0.8
        elif any(word in combined_text for word in ["secondary", "high school", "gcse", "a-level", "year 7", "year 8", "year 9", "year 10", "year 11"]):
            grade_level = "secondary"
            grade_level_confidence = 0.8
        elif any(word in combined_text for word in ["adult", "mature", "university", "college", "professional", "continuing education"]):
            grade_level = "adult-ed"
            grade_level_confidence = 0.7
        
        # Generate suggested tags - following mcp-smart-notes pattern
        suggested_tags = []
        if any(word in combined_text for word in ["reading", "read"]):
            suggested_tags.append("reading")
        if any(word in combined_text for word in ["writing", "write"]):
            suggested_tags.append("writing")
        if any(word in combined_text for word in ["number", "math", "calculation"]):
            suggested_tags.append("numeracy")
        if any(word in combined_text for word in ["comprehension", "understand"]):
            suggested_tags.append("comprehension")
        if any(word in combined_text for word in ["vocabulary", "words"]):
            suggested_tags.append("vocabulary")
        if any(word in combined_text for word in ["creative", "art", "design"]):
            suggested_tags.append("creative")
        
        # Limit to 4 tags
        suggested_tags = suggested_tags[:4]
        
        # Calculate overall confidence
        overall_confidence = (content_type_confidence + subject_confidence + 
                             difficulty_confidence + grade_level_confidence) / 4
        
        return {
            "content_type": content_type,
            "content_type_confidence": content_type_confidence,
            "subject": subject,
            "subject_confidence": subject_confidence,
            "difficulty": difficulty,
            "difficulty_confidence": difficulty_confidence,
            "grade_level": grade_level,
            "grade_level_confidence": grade_level_confidence,
            "suggested_tags": suggested_tags,
            "overall_confidence": overall_confidence,
            "analysis_method": "fallback"
        }
    
    def extract_text_from_file(self, file_path: str, mime_type: Optional[str] = None) -> str:
        """
        Extract text content from various file formats
        Enhanced from original content_analysis.py implementation
        """
        try:
            file_path_obj = Path(file_path)
            file_extension = file_path_obj.suffix.lower()
            
            # Text files
            if file_extension in ['.txt', '.md', '.rst']:
                with open(file_path_obj, 'r', encoding='utf-8') as f:
                    return f.read()
            
            # PDF files
            if file_extension == '.pdf' and PDF_AVAILABLE:
                try:
                    with open(file_path_obj, 'rb') as f:
                        pdf_reader = PyPDF2.PdfReader(f)
                        text = ""
                        for page in pdf_reader.pages:
                            text += page.extract_text() + "\n"
                        return text
                except Exception as e:
                    logging.warning(f"PDF extraction failed: {e}")
            
            # Word documents
            if file_extension in ['.docx', '.doc'] and DOCX_AVAILABLE:
                try:
                    doc = Document(str(file_path_obj))
                    text = ""
                    for paragraph in doc.paragraphs:
                        text += paragraph.text + "\n"
                    return text
                except Exception as e:
                    logging.warning(f"DOCX extraction failed: {e}")
            
            # PowerPoint presentations
            if file_extension in ['.pptx', '.ppt'] and PPTX_AVAILABLE:
                try:
                    prs = Presentation(str(file_path_obj))
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    return text
                except Exception as e:
                    logging.warning(f"PPTX extraction failed: {e}")
            
            # Fallback - try to read as text
            try:
                with open(file_path_obj, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            except Exception:
                pass
                
            return ""
            
        except Exception as e:
            logging.error(f"Content extraction failed for {file_path}: {e}")
            return ""
    
    def analyze_uploaded_content(self, file, metadata: Dict[str, str]) -> Dict[str, Any]:
        """
        Main endpoint function for analyzing uploaded content
        Enhanced version of analyze_content_endpoint from content_analysis.py
        """
        
        # Extract metadata
        title = metadata.get('title', '')
        description = metadata.get('description', '')
        filename = metadata.get('filename', getattr(file, 'filename', ''))
        
        # Save file temporarily for content extraction
        temp_file = None
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, 
                                           suffix=Path(filename).suffix if filename else '.tmp') as tf:
                if hasattr(file, 'save'):
                    file.save(tf.name)
                else:
                    tf.write(file.read())
                temp_file = tf.name
            
            # Extract content from file
            mime_type = getattr(file, 'content_type', None)
            content = self.extract_text_from_file(temp_file, mime_type)
            
            # Combine all text for analysis
            combined_content = f"{title}\n{description}\n{content}".strip()
            
            # Perform analysis using our enhanced method
            analysis = self.analyze_educational_content(
                title=title,
                content=combined_content,
                filename=filename
            )
            
            # Add metadata
            analysis['analysis_timestamp'] = str(Path().cwd())  # Use current time in production
            analysis['has_llm'] = self.client is not None
            analysis['extracted_content_length'] = len(content)
            
            return {
                'status': 'success',
                'analysis': analysis,
                'metadata': {
                    'content_extracted': len(content) > 0,
                    'analysis_method': analysis.get('analysis_method', 'unknown'),
                    'filename': filename,
                    'model_used': analysis.get('model_used', 'fallback')
                }
            }
            
        except Exception as e:
            logging.error(f"Content analysis failed: {e}")
            return {
                'status': 'error',
                'message': str(e)
            }
        
        finally:
            # Clean up temporary file
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except Exception as e:
                    logging.warning(f"Failed to clean up temp file: {e}")
    
    def get_analyzer_status(self) -> Dict[str, Any]:
        """Get status information about the analyzer"""
        return {
            'ollama_available': OLLAMA_AVAILABLE,
            'llm_connected': self.client is not None,
            'model': self.model,
            'supported_formats': {
                'pdf': PDF_AVAILABLE,
                'docx': DOCX_AVAILABLE,
                'pptx': PPTX_AVAILABLE,
                'text': True
            },
            'educational_categories': EDUCATIONAL_CATEGORIES,
            'subject_areas': SUBJECT_AREAS,
            'difficulty_levels': DIFFICULTY_LEVELS,
            'grade_targets': GRADE_TARGETS
        }
    
    def _extract_json_from_response(self, response_text: str) -> Dict[str, Any]:
        """
        Extract JSON from LLM response, handling various formats
        LLMs sometimes include explanatory text before/after JSON
        """
        # First try direct parsing
        try:
            return json.loads(response_text)
        except json.JSONDecodeError:
            pass
        
        # Try to find JSON within the response
        # Look for JSON object pattern
        
        # Pattern 1: Find content between first { and last }
        json_match = re.search(r'\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}', response_text, re.DOTALL)
        if json_match:
            try:
                return json.loads(json_match.group())
            except json.JSONDecodeError:
                pass
        
        # Pattern 2: Look for ```json blocks (fixed regex)
        json_block_match = re.search(r'```(?:json)?\s*(.*?)\s*```', response_text, re.DOTALL)
        if json_block_match:
            try:
                json_content = json_block_match.group(1).strip()
                return json.loads(json_content)
            except json.JSONDecodeError:
                pass
        
        # Pattern 3: Split by newlines and try each line
        lines = response_text.strip().split('\n')
        for line in lines:
            if line.strip().startswith('{') and line.strip().endswith('}'):
                try:
                    return json.loads(line)
                except json.JSONDecodeError:
                    continue
        
        # If all else fails, raise with helpful error
        logging.error(f"Failed to extract JSON from LLM response: {response_text[:500]}...")
        raise json.JSONDecodeError(
            "Could not extract valid JSON from LLM response",
            response_text, 0
        )
    
    def generate_complete_metadata(self, content: str, filename: str) -> Dict[str, Any]:
        """
        Task 2.2 Enhancement: Single LLM call generates ALL metadata
        Zero-touch processing - generates complete database fields
        """
        if not self.client:
            raise Exception("LLM not available - Auto-upload requires Ollama to be running with 'ollama serve'")
        
        # Enhanced prompt for complete metadata generation
        metadata_prompt = f"""Generate complete database metadata for this educational content.

Content: "{content[:2000]}{'...' if len(content) > 2000 else ''}"
Filename: "{filename}"

Return ONLY valid JSON with ALL fields filled:
{{
    "title": "Descriptive title based on content (not filename)",
    "description": "2-3 sentence summary covering learning objectives and key content",
    "subject": "Must be EXACTLY one of: English, Religious Education, Learning Support, Other",
    "content_type": "Must be EXACTLY one of: lesson-plan, worksheet, assessment, resource, activity", 
    "keywords": "Comma-separated search keywords (5-10 relevant terms)",
    "estimated_duration": 30,
    "grade_level": "Must be EXACTLY one of: early-years, primary, secondary, adult-ed",
    "difficulty": "Must be EXACTLY one of: beginner, intermediate, advanced",
    "suggested_tags": ["ONLY use tags from: worksheet, lesson-plan, assessment, interactive, homework, group-work, individual, beginner, advanced, resource, activity"],
    "learning_objectives": "Brief list of what students will learn or achieve",
    "materials_needed": "Any materials or resources required (if applicable)"
}}

Guidelines:
- Title should be clear and descriptive (15-60 characters)
- Description should highlight educational value
- Keywords should include topic terms, skills, and concepts
- Duration in minutes for typical classroom use
- Focus on educational value and practical classroom use
- Make the title engaging and specific to the content
- For suggested_tags: ONLY use tags from the list provided, do NOT create new tags"""

        try:
            # Call LLM for complete metadata generation
            logging.info(f"🤖 Calling LLM for metadata generation with model: {self.model}")
            response = self.client.chat(
                model=self.model,
                messages=[{
                    "role": "user",
                    "content": metadata_prompt
                }]
            )
            
            response_text = (response.message.content or "").strip()
            logging.debug(f"LLM raw response (first 500 chars): {response_text[:500]}...")
            
            try:
                # Use robust JSON extraction method
                metadata = self._extract_json_from_response(response_text)
                
                # Validate and normalize metadata
                metadata = self._validate_and_normalize_metadata(metadata)
                
                # Add generation metadata
                metadata['auto_processed'] = True
                metadata['generation_model'] = self.model
                metadata['categorization_confidence'] = 0.9  # High confidence for complete generation
                
                logging.info(f"✅ Complete metadata generated successfully")
                return {
                    'status': 'success',
                    'metadata': metadata
                }
                
            except json.JSONDecodeError as e:
                logging.error(f"JSON parsing failed: {e}")
                logging.error(f"Full LLM response: {response_text}")
                
                # Try fallback metadata generation with basic analysis
                logging.warning("Attempting fallback metadata generation...")
                fallback = self._generate_fallback_metadata(content, filename)
                if fallback:
                    return {
                        'status': 'success',
                        'metadata': fallback
                    }
                
                raise Exception("LLM returned invalid JSON for metadata generation")
                
        except Exception as e:
            logging.error(f"❌ Metadata generation failed: {e}")
            
            # Last resort: basic metadata
            try:
                logging.warning("Using last resort basic metadata...")
                basic_metadata = self._generate_basic_metadata(content, filename)
                return {
                    'status': 'success',
                    'metadata': basic_metadata
                }
            except:
                raise Exception(f"Metadata generation failed: {str(e)}")
    
    def _validate_and_normalize_metadata(self, metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Validate and normalize generated metadata to ensure all required fields exist"""
        
        # Required fields with defaults
        defaults = {
            'title': 'Untitled Educational Content',
            'description': 'Educational content for classroom use',
            'subject': 'Other',
            'content_type': 'resource',
            'keywords': '',
            'estimated_duration': 30,
            'grade_level': 'primary',
            'difficulty': 'intermediate',
            'suggested_tags': [],
            'learning_objectives': '',
            'materials_needed': ''
        }
        
        # Ensure all fields exist
        for field, default_value in defaults.items():
            if field not in metadata:
                metadata[field] = default_value
        
        # Validate subject (handle comma-separated values)
        subject_value = metadata['subject']
        if isinstance(subject_value, str) and ',' in subject_value:
            # Take the first valid subject from comma-separated list
            subjects = [s.strip() for s in subject_value.split(',')]
            for subj in subjects:
                if subj in SUBJECT_AREAS:
                    metadata['subject'] = subj
                    break
                # Try mapping
                subject_mapping = {
                    'religious education': 'Religious Education',
                    'religious-education': 'Religious Education',
                    'learning support': 'Learning Support',
                    'learning-support': 'Learning Support',
                    'english': 'English',
                    'science': 'Other',  # Map science to Other since it's not in our list
                    'history': 'Other',  # Map history to Other
                    'mathematics': 'Other'  # Map mathematics to Other since it's not in our list
                }
                mapped = subject_mapping.get(subj.lower())
                if mapped and mapped in SUBJECT_AREAS:
                    metadata['subject'] = mapped
                    break
            else:
                metadata['subject'] = 'Other'
        elif metadata['subject'] not in SUBJECT_AREAS:
            # Try to map common variations
            subject_mapping = {
                'religious education': 'Religious Education',
                'religious-education': 'Religious Education',
                'learning support': 'Learning Support',
                'learning-support': 'Learning Support',
                'english': 'English',
                'science': 'Other',  # Map science to Other since it's not in our list
                'mathematics': 'Other'  # Map mathematics to Other since it's not in our list
            }
            mapped = subject_mapping.get(metadata['subject'].lower(), 'Other')
            metadata['subject'] = mapped
        
        # Validate content type
        if metadata['content_type'] not in EDUCATIONAL_CATEGORIES:
            metadata['content_type'] = 'resource'
        
        # Validate difficulty
        if metadata['difficulty'] not in DIFFICULTY_LEVELS:
            metadata['difficulty'] = 'intermediate'
        
        # Validate grade level (handle comma-separated values)
        grade_value = metadata['grade_level']
        if isinstance(grade_value, str) and ',' in grade_value:
            # Take the first valid grade from comma-separated list
            grades = [g.strip() for g in grade_value.split(',')]
            for grade in grades:
                if grade in GRADE_TARGETS:
                    metadata['grade_level'] = grade
                    break
                # Try mapping
                grade_mapping = {
                    'early years': 'early-years',
                    'kindergarten': 'early-years',
                    'elementary': 'primary',
                    'middle school': 'secondary',
                    'high school': 'secondary',
                    'adult': 'adult-ed',
                    'adult education': 'adult-ed'
                }
                mapped = grade_mapping.get(grade.lower())
                if mapped and mapped in GRADE_TARGETS:
                    metadata['grade_level'] = mapped
                    break
            else:
                metadata['grade_level'] = 'primary'  # Default
        elif metadata['grade_level'] not in GRADE_TARGETS:
            # Try to map common variations
            grade_mapping = {
                'early years': 'early-years',
                'kindergarten': 'early-years',
                'elementary': 'primary',
                'middle school': 'secondary',
                'high school': 'secondary',
                'adult': 'adult-ed',
                'adult education': 'adult-ed'
            }
            mapped = grade_mapping.get(metadata['grade_level'].lower(), 'primary')
            metadata['grade_level'] = mapped
        
        # Ensure duration is an integer
        try:
            metadata['estimated_duration'] = int(metadata['estimated_duration'])
        except:
            metadata['estimated_duration'] = 30
        
        # Validate and limit suggested tags
        ALLOWED_TAGS = ['worksheet', 'lesson-plan', 'assessment', 'interactive', 
                        'homework', 'group-work', 'individual', 'beginner', 
                        'advanced', 'resource', 'activity']
        
        if isinstance(metadata['suggested_tags'], list):
            # Filter to only allowed tags
            valid_tags = [tag for tag in metadata['suggested_tags'] if tag in ALLOWED_TAGS]
            metadata['suggested_tags'] = valid_tags[:4]  # Limit to 4 tags
        else:
            metadata['suggested_tags'] = []
        
        # Ensure keywords is a string
        if isinstance(metadata['keywords'], list):
            metadata['keywords'] = ', '.join(metadata['keywords'])
        
        return metadata
    
    def _generate_fallback_metadata(self, content: str, filename: str) -> Dict[str, Any]:
        """Generate metadata using keyword-based analysis when LLM fails"""
        try:
            # Use the existing _fallback_analysis method
            analysis = self._fallback_analysis("", content, filename)
            
            # Extract filename without extension for title
            from pathlib import Path
            base_name = Path(filename).stem.replace('_', ' ').replace('-', ' ')
            
            # Create complete metadata from fallback analysis
            metadata = {
                'title': base_name.title()[:60],  # Clean title from filename
                'description': f"Educational content analyzed from {filename}. Content appears to be a {analysis['content_type'].replace('-', ' ')}.",
                'subject': analysis['subject'],
                'content_type': analysis['content_type'],
                'keywords': ', '.join(analysis.get('suggested_tags', [])),
                'estimated_duration': 30,
                'grade_level': analysis['grade_level'],
                'difficulty': analysis['difficulty'],
                'suggested_tags': analysis.get('suggested_tags', []),
                'learning_objectives': f"Students will work with {analysis['content_type'].replace('-', ' ')} materials",
                'materials_needed': 'Standard classroom supplies',
                'auto_processed': True,
                'generation_model': 'fallback',
                'categorization_confidence': analysis['overall_confidence']
            }
            
            return metadata
            
        except Exception as e:
            logging.error(f"Fallback metadata generation failed: {e}")
            raise
    
    def _generate_basic_metadata(self, content: str, filename: str) -> Dict[str, Any]:
        """Generate very basic metadata as last resort"""
        from pathlib import Path
        base_name = Path(filename).stem.replace('_', ' ').replace('-', ' ')
        
        # Determine basic content type from filename
        filename_lower = filename.lower()
        content_type = 'resource'
        if 'worksheet' in filename_lower:
            content_type = 'worksheet'
        elif 'lesson' in filename_lower or 'plan' in filename_lower:
            content_type = 'lesson-plan'
        elif 'test' in filename_lower or 'quiz' in filename_lower or 'exam' in filename_lower:
            content_type = 'assessment'
        
        return {
            'title': base_name.title()[:60],
            'description': f"Educational {content_type.replace('-', ' ')} uploaded from {filename}",
            'subject': 'Other',
            'content_type': content_type,
            'keywords': content_type.replace('-', ', '),
            'estimated_duration': 30,
            'grade_level': 'primary',
            'difficulty': 'intermediate',
            'suggested_tags': [content_type.split('-')[0]],
            'learning_objectives': 'To be determined',
            'materials_needed': 'Standard classroom supplies',
            'auto_processed': True,
            'generation_model': 'basic',
            'categorization_confidence': 0.3
        }
    
    def auto_process_and_save(self, file, upload_path: str) -> Dict[str, Any]:
        """
        Task 2.2 Enhancement: Zero-touch processing pipeline
        1. Extract text content
        2. Generate metadata via LLM
        3. Return all data ready for database save
        Does NOT save to database - that's handled by the API endpoint
        """
        
        # Get filename for context
        filename = getattr(file, 'filename', 'uploaded_file')
        
        # Save file temporarily for content extraction
        temp_file = None
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, 
                                           suffix=Path(filename).suffix if filename else '.tmp') as tf:
                if hasattr(file, 'save'):
                    file.save(tf.name)
                else:
                    tf.write(file.read())
                temp_file = tf.name
            
            # Extract content from file
            mime_type = getattr(file, 'content_type', None)
            content = self.extract_text_from_file(temp_file, mime_type)
            
            if not content:
                # If no text extracted, use filename as context
                content = f"Educational file: {filename}"
            
            # Generate complete metadata via LLM
            metadata_result = self.generate_complete_metadata(content, filename)
            
            if metadata_result['status'] != 'success':
                raise Exception("Failed to generate metadata")
            
            metadata = metadata_result['metadata']
            
            # Prepare complete data for database save
            auto_data = {
                'title': metadata['title'],
                'description': metadata['description'],
                'subject': metadata['subject'],
                'content_type': metadata['content_type'],
                'keywords': metadata['keywords'],
                'grade_level': metadata['grade_level'],
                'difficulty_level': metadata['difficulty'],
                'duration': metadata['estimated_duration'],
                'suggested_tags': metadata['suggested_tags'],
                'auto_categorized': True,
                'categorization_confidence': metadata.get('categorization_confidence', 0.9),
                'content': content[:5000] if content else '',  # Store first 5000 chars for search
                'original_filename': filename,
                'mime_type': mime_type or 'application/octet-stream',
                'generated_metadata': json.dumps({
                    'learning_objectives': metadata.get('learning_objectives', ''),
                    'materials_needed': metadata.get('materials_needed', ''),
                    'generation_model': metadata.get('generation_model', self.model)
                })
            }
            
            logging.info(f"✅ Auto-processing completed for: {filename}")
            
            return {
                'status': 'success',
                'auto_data': auto_data,
                'metadata': metadata,
                'content_extracted': len(content) > 0,
                'content_length': len(content)
            }
            
        except Exception as e:
            logging.error(f"❌ Auto-processing failed: {e}")
            return {
                'status': 'error',
                'message': str(e)
            }
        
        finally:
            # Clean up temporary file
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except Exception as e:
                    logging.warning(f"Failed to clean up temp file: {e}") 